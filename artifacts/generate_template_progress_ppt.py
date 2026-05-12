#!/usr/bin/env python3
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
TMP_TEMPLATE = ROOT / ".tmp" / "template_blue.pptx"
OUTPUT_TMP = ROOT / ".tmp" / "seu_progress_template_output.pptx"
OUTPUT_FINAL = ROOT / "artifacts" / "毕业进展PPT_东大模板版_v3.pptx"


LAYOUT_LIBRARY = {
    "cover": 0,
    "section": 2,
    "narrative": 3,
    "three_cards": 4,
    "summary_band": 5,
    "two_column_cards": 6,
    "three_panel": 7,
    "four_blocks": 8,
    "image_text": 11,
    "thanks": 14,
}


def rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color.replace("#", "").upper())


def delete_slide(prs: Presentation, slide_index: int) -> None:
    slide_id_list = prs.slides._sldIdLst  # type: ignore[attr-defined]
    slide_ids = list(slide_id_list)
    slide_id = slide_ids[slide_index]
    prs.part.drop_rel(slide_id.rId)
    slide_id_list.remove(slide_id)


def remove_shape(shape) -> None:
    element = shape._element
    parent = element.getparent()
    if parent is not None:
        parent.remove(element)


def iter_text_shapes(shape):
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            yield from iter_text_shapes(child)
    elif getattr(shape, "has_text_frame", False):
        yield shape


def clear_all_text(slide) -> None:
    for shape in slide.shapes:
        for text_shape in iter_text_shapes(shape):
            text_shape.text = ""


def set_text(
    shape,
    lines: list[str] | str,
    *,
    font_size: int = 20,
    bold: bool = False,
    color: str = "17324D",
    align: PP_ALIGN = PP_ALIGN.LEFT,
    font_name: str = "微软雅黑",
    line_spacing: float | None = 1.1,
) -> None:
    if not getattr(shape, "has_text_frame", False):
        return
    if isinstance(lines, str):
        lines = [lines]
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    for idx, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        paragraph.alignment = align
        if line_spacing is not None:
            paragraph.line_spacing = line_spacing
        run = paragraph.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)


def add_box(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    lines: list[str] | str,
    *,
    font_size: int = 18,
    bold: bool = False,
    color: str = "17324D",
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> None:
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    set_text(
        shape,
        lines,
        font_size=font_size,
        bold=bold,
        color=color,
        align=align,
    )


def render_cover(slide) -> None:
    clear_all_text(slide)
    for idx in [7, 8]:
        remove_shape(slide.shapes[idx])
    set_text(slide.shapes[4], "毕业设计进展汇报", font_size=28, bold=True, color="FFFFFF")
    set_text(slide.shapes[6], "基于轻量化联邦学习的物联网恶意流量检测系统", font_size=13, color="DCE7F1")
    add_box(slide, 1.62, 4.86, 3.3, 0.34, "汇报人：朱世豪", font_size=12, color="EAF2F8")
    add_box(slide, 1.62, 5.26, 3.5, 0.34, "指导教师：赵瑞杰", font_size=12, color="EAF2F8")
    add_box(slide, 1.62, 5.66, 3.3, 0.34, "时间：2026年5月", font_size=12, color="EAF2F8")


def render_section(slide) -> None:
    clear_all_text(slide)
    add_box(slide, 6.28, 2.75, 5.1, 0.65, "研究背景与课题定位", font_size=24, bold=True, color="FFFFFF")
    set_text(slide.shapes[6], "围绕隐私保护、通信开销与异构协同展开", font_size=13, color="DCE7F1")


def render_narrative_background(slide) -> None:
    clear_all_text(slide)
    set_text(slide.shapes[0], "课题背景与研究目标", font_size=24, bold=True)
    set_text(slide.shapes[5], "研究主线", font_size=12, bold=True, color="FFFFFF", align=PP_ALIGN.CENTER)
    set_text(
        slide.shapes[7],
        [
            "随着物联网设备规模快速增长，恶意流量检测面临数据难集中共享、终端算力受限和客户端分布异构等问题。",
            "传统集中式训练存在隐私泄露风险，FedAvg 在 Non-IID 场景下又容易出现收敛不稳定和类别偏置。",
            "本课题以 FedProto 为核心协同机制，尝试通过原型通信替代参数通信，降低联邦训练开销。"
        ],
        font_size=15,
        color="4A5F75",
    )
    set_text(slide.shapes[10], "本课题拟解决的问题", font_size=13, bold=True, color="17324D", align=PP_ALIGN.CENTER)
    set_text(
        slide.shapes[11],
        [
            "在原始数据不出本地的条件下实现协同检测",
            "通过低维类别原型降低通信负担",
            "支持 MLP、CNN1D 等异构客户端共同训练"
        ],
        font_size=12,
        color="4A5F75",
    )


def render_three_cards_progress(slide) -> None:
    clear_all_text(slide)
    set_text(slide.shapes[0], "目前已完成的主要工作", font_size=24, bold=True)
    set_text(slide.shapes[16], "研究方案", font_size=12, bold=True, color="17324D", align=PP_ALIGN.CENTER)
    set_text(
        slide.shapes[17],
        [
            "完成开题报告、论文绪论与相关技术部分撰写",
            "明确“问题分析-方案设计-系统实现-实验验证”研究路线"
        ],
        font_size=11,
        color="4A5F75",
    )
    set_text(slide.shapes[11], "系统实现", font_size=12, bold=True, color="17324D", align=PP_ALIGN.CENTER)
    set_text(
        slide.shapes[12],
        [
            "完成数据预处理、训练入口、结果归档与可视化界面联动",
            "实现 MLP、CNN1D、Transformer1D 三类本地模型"
        ],
        font_size=11,
        color="4A5F75",
    )
    set_text(slide.shapes[17], "实验进展", font_size=12, bold=True, color="17324D", align=PP_ALIGN.CENTER)
    set_text(
        slide.shapes[18],
        [
            "完成 Local、FedAvg、FedProto 同构实验",
            "完成 MLP + CNN1D 异构 FedProto 扩展实验"
        ],
        font_size=11,
        color="4A5F75",
    )


def render_three_panel_system(slide) -> None:
    clear_all_text(slide)
    set_text(slide.shapes[2], "系统设计与实现进展", font_size=24, bold=True)
    set_text(slide.shapes[8], "数据与训练链路", font_size=12, bold=True, color="17324D", align=PP_ALIGN.CENTER)
    set_text(slide.shapes[12], "模型与协同机制", font_size=12, bold=True, color="17324D", align=PP_ALIGN.CENTER)
    set_text(slide.shapes[15], "界面与管理能力", font_size=12, bold=True, color="17324D", align=PP_ALIGN.CENTER)
    set_text(
        slide.shapes[9],
        [
            "完成 77 维统计特征清洗、编码、归一化与 Non-IID 客户端划分。",
            "统一调度 Local、FedAvg、FedProto 三类训练流程。"
        ],
        font_size=11,
        color="4A5F75",
    )
    set_text(
        slide.shapes[11],
        [
            "实现统一特征输出维度，支持原型提取与异构协同。",
            "FedProto 已打通客户端原型提取与服务端聚合流程。"
        ],
        font_size=11,
        color="4A5F75",
    )
    set_text(
        slide.shapes[14],
        [
            "Electron 桌面端已能组织实验配置、日志与结果展示。",
            "当前已形成训练、归档、复测的基本闭环。"
        ],
        font_size=11,
        color="4A5F75",
    )


def render_two_column_setup(slide) -> None:
    clear_all_text(slide)
    set_text(slide.shapes[0], "实验设置与数据组织", font_size=24, bold=True)
    set_text(slide.shapes[5], "实验配置", font_size=12, bold=True, color="17324D", align=PP_ALIGN.CENTER)
    set_text(slide.shapes[6], "数据组织", font_size=12, bold=True, color="17324D", align=PP_ALIGN.CENTER)
    set_text(
        slide.shapes[12],
        [
            "20 个客户端、15 类任务、77 维输入特征、特征维度 512",
            "对比算法：Local、FedAvg、FedProto",
            "同构模型：MLP、CNN1D、Transformer1D"
        ],
        font_size=11,
        color="4A5F75",
    )
    set_text(
        slide.shapes[14],
        [
            "每个客户端保留 5 个类别，训练集 375 条、测试集 125 条",
            "全局共 7500 条训练样本、2500 条测试样本",
            "已构建符合物联网场景特征的 Non-IID 数据划分"
        ],
        font_size=11,
        color="4A5F75",
    )
    set_text(slide.shapes[8], "项目现状", font_size=12, bold=True, color="17324D", align=PP_ALIGN.CENTER)
    set_text(
        slide.shapes[11],
        [
            "项目代码、实验结果、界面原型和论文草稿已经基本对齐。"
        ],
        font_size=11,
        color="4A5F75",
    )
    set_text(slide.shapes[10], "阶段判断", font_size=12, bold=True, color="17324D", align=PP_ALIGN.CENTER)
    set_text(
        slide.shapes[9],
        [
            "当前已从“方案搭建”进入“结果收口与表达优化”阶段。"
        ],
        font_size=11,
        color="4A5F75",
    )


def render_image_text_results(slide) -> None:
    clear_all_text(slide)
    set_text(slide.shapes[1], "阶段性实验结果", font_size=24, bold=True)
    set_text(slide.shapes[8], "结果图示", font_size=12, bold=True, color="17324D", align=PP_ALIGN.CENTER)
    set_text(slide.shapes[14], "关键结论", font_size=12, bold=True, color="17324D", align=PP_ALIGN.CENTER)
    left_pic = ROOT / ".tmp" / "chapter5_word_assets" / "fig5-1-accuracy-f1.png"
    if left_pic.exists():
        slide.shapes.add_picture(str(left_pic), slide.shapes[5].left, slide.shapes[5].top, slide.shapes[5].width, slide.shapes[5].height)
    set_text(
        slide.shapes[10],
        [
            "CNN1D + FedProto 取得当前最优结果：Acc = 0.9484，F1 = 0.9481。",
            "MLP + FedProto 与本地训练结果接近，说明原型对齐未显著削弱轻量模型性能。"
        ],
        font_size=12,
        color="4A5F75",
    )
    set_text(slide.shapes[12], "异构扩展", font_size=12, bold=True, color="17324D", align=PP_ALIGN.CENTER)
    set_text(
        slide.shapes[13],
        [
            "MLP + CNN1D 异构 FedProto 达到 Acc = 0.9416。",
            "最优结果出现在第 832 轮，验证了统一原型空间设计可行。"
        ],
        font_size=11,
        color="4A5F75",
    )


def render_four_blocks_issue_plan(slide) -> None:
    clear_all_text(slide)
    set_text(slide.shapes[2], "当前问题与下一步计划", font_size=24, bold=True)
    set_text(slide.shapes[8], "当前不足", font_size=12, bold=True, color="17324D")
    set_text(
        slide.shapes[11],
        [
            "实验主要基于离线数据集，尚未覆盖真实边缘设备环境。"
        ],
        font_size=10,
        color="4A5F75",
    )
    set_text(slide.shapes[12], "问题二", font_size=12, bold=True, color="17324D")
    set_text(
        slide.shapes[13],
        [
            "异构实验目前主要验证 MLP 与 CNN1D 组合，范围仍有限。"
        ],
        font_size=10,
        color="4A5F75",
    )
    add_box(slide, 0.95, 5.20, 5.15, 0.34, "下一步一", font_size=12, bold=True, color="17324D")
    add_box(
        slide,
        0.95,
        5.58,
        5.20,
        0.90,
        ["继续完善论文正文、图表排版和实验分析表达。"],
        font_size=10,
        color="4A5F75",
    )
    add_box(slide, 6.85, 5.20, 5.15, 0.34, "下一步二", font_size=12, bold=True, color="17324D")
    add_box(
        slide,
        6.85,
        5.58,
        5.20,
        0.90,
        ["补充系统界面截图、关键流程图与最终汇报材料，并结合老师意见调整结论表述。"],
        font_size=10,
        color="4A5F75",
    )


def render_thanks(slide) -> None:
    clear_all_text(slide)
    for idx in [7, 8]:
        remove_shape(slide.shapes[idx])
    set_text(slide.shapes[4], "感谢各位老师聆听", font_size=28, bold=True, color="FFFFFF")
    set_text(slide.shapes[6], "恳请各位老师批评指正", font_size=13, color="DCE7F1")
    add_box(slide, 1.62, 4.86, 4.0, 0.34, "项目主题：FedProto + 轻量化检测", font_size=12, color="EAF2F8")
    add_box(slide, 1.62, 5.26, 4.7, 0.34, "汇报材料基于论文与开题报告整理", font_size=12, color="EAF2F8")


CONTENT_PLAN = [
    ("cover", render_cover),
    ("section", render_section),
    ("narrative", render_narrative_background),
    ("three_cards", render_three_cards_progress),
    ("two_column_cards", render_two_column_setup),
    ("three_panel", render_three_panel_system),
    ("four_blocks", render_four_blocks_issue_plan),
    ("image_text", render_image_text_results),
    ("thanks", render_thanks),
]


def main() -> None:
    prs = Presentation(str(TMP_TEMPLATE))
    keep = sorted(LAYOUT_LIBRARY[key] for key, _ in CONTENT_PLAN)
    for idx in range(len(prs.slides) - 1, -1, -1):
        if idx not in keep:
            delete_slide(prs, idx)

    for slide, (layout_key, renderer) in zip(prs.slides, CONTENT_PLAN):
        renderer(slide)

    prs.save(str(OUTPUT_TMP))
    OUTPUT_FINAL.write_bytes(OUTPUT_TMP.read_bytes())
    print(OUTPUT_FINAL)


if __name__ == "__main__":
    main()
