#!/usr/bin/env python3
import argparse
from pathlib import Path

from pptx import Presentation


def main():
    parser = argparse.ArgumentParser(description="Extract visible text from a PPTX for quick QA.")
    parser.add_argument("pptx_path", help="Path to the .pptx file to inspect.")
    args = parser.parse_args()

    prs = Presentation(str(Path(args.pptx_path)))
    print(f"slides: {len(prs.slides)}")
    for slide_index, slide in enumerate(prs.slides, start=1):
        print(f"\n# slide {slide_index}")
        seen = False
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            text = shape.text.strip()
            if not text:
                continue
            seen = True
            print(text)
        if not seen:
            print("[no visible text]")


if __name__ == "__main__":
    main()
