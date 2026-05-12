#!/usr/bin/env python3
import argparse
import json
from pathlib import Path
import re

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


DEFAULT_THEME = {
    "primary": "1F3A5F",
    "secondary": "4C6A92",
    "accent": "E07A5F",
    "light": "E8F1F8",
    "bg": "F7FAFC",
    "font_cn": "Microsoft YaHei",
    "font_en": "Arial",
}


def has_cjk(text: str) -> bool:
    return bool(re.search(r"[\u3400-\u4dbf\u4e00-\u9fff\uf900-\ufaff]", text))


def pick_font(text: str, theme: dict) -> str:
    return theme["font_cn"] if has_cjk(text) else theme["font_en"]


def rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color.strip().replace("#", "").upper())


def merge_theme(data: dict) -> dict:
    theme = dict(DEFAULT_THEME)
    theme.update(data.get("theme", {}))
    return theme


def set_run_style(run, text: str, theme: dict, size: int, color: str, bold: bool = False):
    run.text = text
    run.font.name = pick_font(text, theme)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    theme,
    size=20,
    color=None,
    bold=False,
    align=PP_ALIGN.LEFT,
    fill=None,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    if fill is not None:
        box.fill.solid()
        box.fill.fore_color.rgb = rgb(fill)
        box.line.fill.background()
    text_frame = box.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.TOP
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    set_run_style(run, text, theme, size=size, color=color or theme["primary"], bold=bold)
    return box


def add_bullets(slide, left, top, width, height, bullets, theme, size=20, color=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = box.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.margin_left = Pt(6)
    text_frame.margin_right = Pt(4)
    text_frame.margin_top = Pt(2)
    text_frame.margin_bottom = Pt(2)

    for index, item in enumerate(bullets):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.level = 0
        paragraph.bullet = True
        run = paragraph.add_run()
        set_run_style(run, str(item), theme, size=size, color=color or theme["secondary"])
    return box


def add_page_number(slide, number: int, theme: dict):
    badge = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL,
        Inches(9.15),
        Inches(5.0),
        Inches(0.42),
        Inches(0.42),
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = rgb(theme["accent"])
    badge.line.fill.background()

    text_frame = badge.text_frame
    text_frame.clear()
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    set_run_style(run, str(number), theme, size=12, color="FFFFFF", bold=True)


def add_title_slide(prs: Presentation, slide_data: dict, theme: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(theme["bg"])

    accent_bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.55),
        Inches(0.85),
        Inches(1.0),
        Inches(3.8),
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = rgb(theme["accent"])
    accent_bar.line.fill.background()

    add_textbox(
        slide,
        Inches(1.75),
        Inches(1.2),
        Inches(7.4),
        Inches(1.4),
        slide_data["title"],
        theme,
        size=28,
        color=theme["primary"],
        bold=True,
    )
    subtitle = slide_data.get("subtitle")
    if subtitle:
        add_textbox(
            slide,
            Inches(1.8),
            Inches(2.55),
            Inches(6.8),
            Inches(1.0),
            subtitle,
            theme,
            size=18,
            color=theme["secondary"],
        )
    footer = slide_data.get("footer")
    if footer:
        add_textbox(
            slide,
            Inches(1.8),
            Inches(4.7),
            Inches(4.8),
            Inches(0.5),
            footer,
            theme,
            size=12,
            color=theme["secondary"],
        )


def add_section_slide(prs: Presentation, slide_data: dict, theme: dict, page_number: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(theme["primary"])

    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.8),
        Inches(1.2),
        Inches(8.0),
        Inches(3.2),
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = rgb(theme["light"])
    panel.line.fill.background()

    add_textbox(
        slide,
        Inches(1.25),
        Inches(1.85),
        Inches(7.0),
        Inches(0.8),
        slide_data["title"],
        theme,
        size=26,
        color=theme["primary"],
        bold=True,
    )
    subtitle = slide_data.get("subtitle")
    if subtitle:
        add_textbox(
            slide,
            Inches(1.25),
            Inches(2.7),
            Inches(6.8),
            Inches(0.8),
            subtitle,
            theme,
            size=16,
            color=theme["secondary"],
        )
    add_page_number(slide, page_number, theme)


def add_bullet_slide(prs: Presentation, slide_data: dict, theme: dict, page_number: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(theme["bg"])

    add_textbox(
        slide,
        Inches(0.7),
        Inches(0.55),
        Inches(8.0),
        Inches(0.7),
        slide_data["title"],
        theme,
        size=24,
        color=theme["primary"],
        bold=True,
    )

    if slide_data.get("lead"):
        add_textbox(
            slide,
            Inches(0.75),
            Inches(1.25),
            Inches(8.4),
            Inches(0.6),
            slide_data["lead"],
            theme,
            size=14,
            color=theme["secondary"],
        )

    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.7),
        Inches(1.7),
        Inches(8.2),
        Inches(2.9),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = rgb("FFFFFF")
    card.line.color.rgb = rgb(theme["light"])

    add_bullets(
        slide,
        Inches(0.95),
        Inches(2.0),
        Inches(7.6),
        Inches(2.3),
        slide_data.get("bullets", []),
        theme,
        size=18,
    )

    if slide_data.get("highlight"):
        highlight = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(6.3),
            Inches(4.45),
            Inches(2.4),
            Inches(0.55),
        )
        highlight.fill.solid()
        highlight.fill.fore_color.rgb = rgb(theme["accent"])
        highlight.line.fill.background()
        add_textbox(
            slide,
            Inches(6.45),
            Inches(4.55),
            Inches(2.05),
            Inches(0.3),
            slide_data["highlight"],
            theme,
            size=11,
            color="FFFFFF",
            bold=True,
            align=PP_ALIGN.CENTER,
        )

    add_page_number(slide, page_number, theme)


def add_two_column_slide(prs: Presentation, slide_data: dict, theme: dict, page_number: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(theme["bg"])

    add_textbox(
        slide,
        Inches(0.7),
        Inches(0.55),
        Inches(8.0),
        Inches(0.7),
        slide_data["title"],
        theme,
        size=24,
        color=theme["primary"],
        bold=True,
    )

    left = slide_data.get("left", {})
    right = slide_data.get("right", {})

    for box_left, content in ((0.7, left), (5.0, right)):
        panel = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(box_left),
            Inches(1.45),
            Inches(3.55),
            Inches(3.3),
        )
        panel.fill.solid()
        panel.fill.fore_color.rgb = rgb("FFFFFF")
        panel.line.color.rgb = rgb(theme["light"])

        heading = content.get("heading", "")
        if heading:
            add_textbox(
                slide,
                Inches(box_left + 0.2),
                Inches(1.72),
                Inches(2.9),
                Inches(0.45),
                heading,
                theme,
                size=16,
                color=theme["primary"],
                bold=True,
            )
        add_bullets(
            slide,
            Inches(box_left + 0.18),
            Inches(2.15),
            Inches(3.05),
            Inches(2.1),
            content.get("bullets", []),
            theme,
            size=15,
        )

    add_page_number(slide, page_number, theme)


def add_summary_slide(prs: Presentation, slide_data: dict, theme: dict, page_number: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(theme["bg"])

    add_textbox(
        slide,
        Inches(0.7),
        Inches(0.55),
        Inches(8.0),
        Inches(0.7),
        slide_data["title"],
        theme,
        size=24,
        color=theme["primary"],
        bold=True,
    )

    add_bullets(
        slide,
        Inches(0.8),
        Inches(1.35),
        Inches(5.2),
        Inches(3.2),
        slide_data.get("bullets", []),
        theme,
        size=18,
    )

    quote = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(6.2),
        Inches(1.6),
        Inches(2.6),
        Inches(2.5),
    )
    quote.fill.solid()
    quote.fill.fore_color.rgb = rgb(theme["primary"])
    quote.line.fill.background()

    add_textbox(
        slide,
        Inches(6.45),
        Inches(2.0),
        Inches(2.1),
        Inches(1.0),
        slide_data.get("highlight", "Key takeaway"),
        theme,
        size=18,
        color="FFFFFF",
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_page_number(slide, page_number, theme)


def normalize_deck(raw: dict) -> dict:
    slides = list(raw.get("slides", []))
    if not slides:
        slides = [
            {
                "type": "title",
                "title": raw.get("title", "Presentation Title"),
                "subtitle": raw.get("subtitle", ""),
                "footer": raw.get("author", ""),
            }
        ]
    elif slides[0].get("type") != "title":
        slides.insert(
            0,
            {
                "type": "title",
                "title": raw.get("title", slides[0].get("title", "Presentation Title")),
                "subtitle": raw.get("subtitle", ""),
                "footer": raw.get("author", ""),
            },
        )
    return {"theme": merge_theme(raw), "slides": slides}


def build_presentation(deck: dict, output_path: Path):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    page_number = 0
    for slide_data in deck["slides"]:
        slide_type = slide_data.get("type", "bullets")
        if slide_type == "title":
            add_title_slide(prs, slide_data, deck["theme"])
            continue

        page_number += 1
        if slide_type == "section":
            add_section_slide(prs, slide_data, deck["theme"], page_number)
        elif slide_type == "two_column":
            add_two_column_slide(prs, slide_data, deck["theme"], page_number)
        elif slide_type == "summary":
            add_summary_slide(prs, slide_data, deck["theme"], page_number)
        else:
            add_bullet_slide(prs, slide_data, deck["theme"], page_number)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))


def main():
    parser = argparse.ArgumentParser(description="Generate a PowerPoint presentation from deck JSON.")
    parser.add_argument("--input", required=True, help="Path to a deck JSON file.")
    parser.add_argument("--output", required=True, help="Path to output .pptx file.")
    args = parser.parse_args()

    deck_path = Path(args.input)
    output_path = Path(args.output)
    data = json.loads(deck_path.read_text(encoding="utf-8"))
    deck = normalize_deck(data)
    build_presentation(deck, output_path)
    print(output_path)


if __name__ == "__main__":
    main()
